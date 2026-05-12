"""
Erstellt KI_Assistent_Architektur.pptx
Erklaert den technischen Ablauf (C4 Level 2) auf Deutsch.

Ausfuehren:
    pip install python-pptx
    python make_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

YELLOW  = RGBColor(0xFF, 0xD0, 0x00)
DARK    = RGBColor(0x1A, 0x1F, 0x2E)
STEEL   = RGBColor(0x3A, 0x4A, 0x6B)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x6B, 0x7A, 0x99)
SUCCESS = RGBColor(0x2E, 0xA0, 0x6A)
WARNING = RGBColor(0xF5, 0xA6, 0x23)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def arrow_right(slide, l, t, w=Inches(0.5), color=MUTED):
    rect(slide, l, t - Inches(0.04), w * 0.75, Inches(0.08), fill=color)
    arrow = slide.shapes.add_shape(13, l, t - Inches(0.15), w, Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def box(slide, l, t, w, h, title, subtitle='', fill=STEEL):
    rect(slide, l, t, w, h, fill=fill, line=YELLOW, line_w=Pt(1.5))
    txt(slide, title, l + Inches(0.12), t + Inches(0.1),
        w - Inches(0.24), Inches(0.35), size=13, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, l + Inches(0.12), t + Inches(0.42),
            w - Inches(0.24), h - Inches(0.5), size=10, color=MUTED,
            align=PP_ALIGN.LEFT)


def header(slide, title, sub=''):
    rect(slide, 0, 0, W, Inches(1.1), fill=DARK)
    txt(slide, title, Inches(0.5), Inches(0.12), Inches(12), Inches(0.58),
        size=26, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
    if sub:
        txt(slide, sub, Inches(0.5), Inches(0.68), Inches(12), Inches(0.35),
            size=13, color=MUTED, align=PP_ALIGN.LEFT)
    rect(slide, 0, Inches(1.1), W, Pt(3), fill=YELLOW)


def footer(slide, note=''):
    rect(slide, 0, H - Inches(0.42), W, Inches(0.42), fill=DARK)
    label = 'KI Assistent  |  LR 1104.03.08  |  Technische Architektur'
    if note:
        label += '  |  ' + note
    txt(slide, label, Inches(0.4), H - Inches(0.36), Inches(12), Inches(0.3),
        size=9, color=MUTED, align=PP_ALIGN.LEFT)


# -- Folie 1: Titel --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=DARK)
rect(sl, 0, H - Inches(1.4), W, Inches(1.4), fill=STEEL)
rect(sl, 0, 0, Inches(0.35), H, fill=YELLOW)
txt(sl, 'KI Assistent', Inches(0.7), Inches(1.8), Inches(11), Inches(1.0),
    size=48, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
txt(sl, 'Technische Architektur & Ablauf',
    Inches(0.7), Inches(2.85), Inches(11), Inches(0.65),
    size=28, color=WHITE, align=PP_ALIGN.LEFT)
txt(sl, 'Liebherr LR 1104.03.08  |  MTP-Demo',
    Inches(0.7), Inches(3.55), Inches(11), Inches(0.5),
    size=16, color=MUTED, align=PP_ALIGN.LEFT)
txt(sl, ('Wie beantwortet der Assistent eine Frage aus dem Betriebshandbuch?\n'
         'Welche Komponenten sind beteiligt -- und warum?'),
    Inches(0.7), Inches(4.3), Inches(10), Inches(0.9),
    size=14, color=RGBColor(0xC0, 0xCB, 0xE0), align=PP_ALIGN.LEFT)
footer(sl)


# -- Folie 2: Systemuebersicht --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LIGHT)
header(sl, 'Systemuebersicht -- alle Komponenten auf einen Blick',
       'C4 Level 2  |  Container-Diagramm')

box(sl, Inches(0.3), Inches(3.0), Inches(1.8), Inches(1.4),
    'Bediener / Techniker', 'Browser oder iPad\nSafari / Chrome', fill=STEEL)

al, at, aw, ah = Inches(2.5), Inches(1.4), Inches(7.5), Inches(4.8)
rect(sl, al, at, aw, ah, fill=RGBColor(0xE8, 0xEE, 0xF8), line=STEEL, line_w=Pt(1.5))
txt(sl, 'KI Assistent  --  Railway (Cloud)',
    al + Inches(0.15), at + Inches(0.08), Inches(5), Inches(0.35),
    size=11, bold=True, color=STEEL, align=PP_ALIGN.LEFT)

bw, bh, bt = Inches(1.55), Inches(1.1), at + Inches(0.55)
box(sl, al+Inches(0.2),  bt, bw, bh, 'Frontend',
    'HTML/CSS/JS\nEin-Datei-App\nWeb Speech API', fill=DARK)
box(sl, al+Inches(2.05), bt, bw, bh, 'Backend API',
    'FastAPI / Python\nPOST /ask\nPOST /errorcode', fill=DARK)
box(sl, al+Inches(3.9),  bt, bw, bh, 'Suchmaschine',
    'rapidfuzz\nIn-Memory-Index\nFuzzy + Keyword', fill=DARK)
box(sl, al+Inches(5.75), bt, bw, bh, 'Stoerungsdatenbank',
    'JSON im Speicher\nStorungscodes\nBeim Start geladen', fill=DARK)

bt2 = at + Inches(2.1)
box(sl, al+Inches(0.2), bt2, Inches(3.4), Inches(0.95),
    'Manual-Seiten  (manuals/*.html)',
    '2.180 HTML-Seiten aus DITA-Export  |  Statische Dateien',
    fill=RGBColor(0x2A, 0x35, 0x55))
box(sl, al+Inches(3.9), bt2, Inches(3.4), Inches(0.95),
    'Such-Indizes  (data/)',
    'content_index.json  |  toc_index.json  |  metadata_index.json',
    fill=RGBColor(0x2A, 0x35, 0x55))

rect(sl, al+Inches(0.2), at+Inches(3.45), Inches(7.1), Inches(0.88),
     fill=RGBColor(0xD0, 0xD8, 0xEC), line=MUTED, line_w=Pt(1))
txt(sl, ('Vorverarbeitungs-CLI  (einmalig nach Manual-Update)\n'
         'parse_toc / parse_metadata / extract_content / extract_errorcodes  ->  schreibt data/ neu'),
    al+Inches(0.35), at+Inches(3.52), Inches(6.8), Inches(0.75),
    size=10, color=STEEL, align=PP_ALIGN.LEFT)

box(sl, Inches(10.6), Inches(2.8), Inches(2.4), Inches(1.5),
    'Anthropic Claude API',
    'claude-haiku-4-5-20251001\nAntwort + Verifikation\nPrompt-Caching aktiv',
    fill=RGBColor(0x6B, 0x3F, 0xC8))

arrow_right(sl, Inches(2.1),  Inches(3.5),  Inches(0.45), color=YELLOW)
arrow_right(sl, Inches(9.95), Inches(3.1),  Inches(0.65), color=YELLOW)
footer(sl, 'Folie 2 / 7')


# -- Folie 3: Ablauf POST /ask --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=DARK)
header(sl, 'Ablauf: Bediener stellt eine Frage',
       'POST /ask  --  ca. 1-2 Sekunden Gesamtzeit')

steps = [
    ('1', 'Bediener tippt oder spricht',
     'Text eingeben oder Mikrofon nutzen (Web Speech API). Audio bleibt auf dem Geraet.',
     YELLOW, DARK),
    ('2', 'Frontend sendet POST /ask',
     'Browser sendet {"question": "..."} als JSON an Railway. Keine Seite wird neu geladen.',
     YELLOW, DARK),
    ('3', 'Suchmaschine findet passende Seiten',
     'rapidfuzz durchsucht In-Memory-Index (2.180 Seiten). Ergebnis: Top-5 Seiten (~5 ms).',
     YELLOW, DARK),
    ('4', 'Claude generiert Antwort  (API-Call 1)',
     'Nur 5 Seiten (nicht das ganze Manual!) an Claude. Strukturierter Text (~800-1.500 ms).',
     RGBColor(0x6B, 0x3F, 0xC8), WHITE),
    ('5', 'Claude prueft Antwort  (API-Call 2)',
     'Kontext aus Cache + Antwort -> BELEGT / TEILWEISE / NICHT_BELEGT (10 Token, ~200-400 ms).',
     RGBColor(0x6B, 0x3F, 0xC8), WHITE),
    ('6', 'Ergebnis wird angezeigt',
     'Antwort, Grounding-Badge und Quellenlinks. Bei NICHT_BELEGT: Fallback statt erfundener Antwort.',
     SUCCESS, DARK),
]
row_h = Inches(0.82)
for i, (num, title, detail, bf, nc) in enumerate(steps):
    t = Inches(1.28) + i * row_h
    rect(sl, Inches(0.25), t+Inches(0.12), Inches(0.42), Inches(0.42), fill=bf)
    txt(sl, num, Inches(0.25), t+Inches(0.1), Inches(0.42), Inches(0.42),
        size=14, bold=True, color=nc, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(0.85), t+Inches(0.04), Inches(4.0), Inches(0.38),
        size=13, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
    txt(sl, detail, Inches(0.85), t+Inches(0.4), Inches(12.2), Inches(0.42),
        size=10, color=RGBColor(0xC0, 0xCB, 0xE0), align=PP_ALIGN.LEFT)
    if i < 5:
        rect(sl, Inches(0.25), t+row_h-Pt(1), Inches(12.8), Pt(1), fill=STEEL)
footer(sl, 'Folie 3 / 7')


# -- Folie 4: Zwei-Call-Design --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LIGHT)
header(sl, 'Warum zwei Claude-Aufrufe?',
       'Antwort-Generierung + unabhaengige Verifikation')

lc, lw = Inches(0.4), Inches(5.8)
rect(sl, lc, Inches(1.3), lw, Inches(5.4),
     fill=RGBColor(0xEE, 0xEC, 0xF9), line=RGBColor(0x6B, 0x3F, 0xC8), line_w=Pt(2))
txt(sl, 'Aufruf 1  --  Antwort generieren',
    lc+Inches(0.2), Inches(1.38), lw-Inches(0.4), Inches(0.45),
    size=14, bold=True, color=RGBColor(0x6B, 0x3F, 0xC8), align=PP_ALIGN.LEFT)

call1 = [
    ('Modell',      'claude-haiku-4-5-20251001'),
    ('Max. Token',  '1.024  (vollstaendige Antwort)'),
    ('System',      'Du bist Assistent fuer Liebherr-Maschinenfuehrer ...  [gecacht]'),
    ('Kontext',     'Top-5 Manual-Seiten als Textblock  [gecacht]\nca. 2.000-6.000 Token'),
    ('Frage',       'Nutzerfrage als letzter Block  [nicht gecacht]'),
    ('Antwort',     'Strukturierter Text mit Schritten & Warnhinweisen'),
    ('Latenz',      '~800-1.500 ms'),
]
ty = Inches(1.92)
for label, value in call1:
    txt(sl, label+':', lc+Inches(0.2), ty, Inches(1.4), Inches(0.35),
        size=10, bold=True, color=STEEL, align=PP_ALIGN.LEFT)
    txt(sl, value, lc+Inches(1.6), ty, lw-Inches(1.8), Inches(0.38),
        size=10, color=DARK, align=PP_ALIGN.LEFT)
    ty += Inches(0.52)

rc, rw = Inches(6.8), Inches(5.8)
rect(sl, rc, Inches(1.3), rw, Inches(5.4),
     fill=RGBColor(0xEE, 0xF8, 0xF1), line=SUCCESS, line_w=Pt(2))
txt(sl, 'Aufruf 2  --  Antwort verifizieren',
    rc+Inches(0.2), Inches(1.38), rw-Inches(0.4), Inches(0.45),
    size=14, bold=True, color=SUCCESS, align=PP_ALIGN.LEFT)

call2 = [
    ('Modell',      'claude-haiku-4-5-20251001'),
    ('Max. Token',  '10  (nur ein Wort erwartet)'),
    ('System',      'Du pruefst ob die Antwort im Kontext belegt ist ...  [gecacht]'),
    ('Kontext',     'Identisch mit Aufruf 1  ->  Cache-Treffer\nKeine erneuten Token-Kosten'),
    ('Input',       'Generierte Antwort aus Aufruf 1'),
    ('Ausgabe',     'BELEGT  /  TEILWEISE  /  NICHT_BELEGT'),
    ('Latenz',      '~200-400 ms  (Kontext aus Cache)'),
]
ty = Inches(1.92)
for label, value in call2:
    txt(sl, label+':', rc+Inches(0.2), ty, Inches(1.4), Inches(0.35),
        size=10, bold=True, color=STEEL, align=PP_ALIGN.LEFT)
    txt(sl, value, rc+Inches(1.6), ty, rw-Inches(1.8), Inches(0.38),
        size=10, color=DARK, align=PP_ALIGN.LEFT)
    ty += Inches(0.52)

rect(sl, Inches(0.4), Inches(6.52), Inches(12.2), Inches(0.62),
     fill=RGBColor(0xFF, 0xF3, 0xCC), line=WARNING, line_w=Pt(1.5))
txt(sl, ('Sicherheitsprinzip: Bei NICHT_BELEGT wird die Antwort verworfen '
         'und durch einen Hinweis ersetzt. Der Assistent erfindet keine Fakten.'),
    Inches(0.6), Inches(6.57), Inches(12.0), Inches(0.5),
    size=11, color=DARK, align=PP_ALIGN.LEFT)
footer(sl, 'Folie 4 / 7')


# -- Folie 5: Prompt-Caching --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=DARK)
header(sl, 'Prompt-Caching -- Kosten & Geschwindigkeit optimieren',
       'Anthropic speichert identische Textbloecke zwischen Anfragen')

blocks = [
    ('System-Prompt\n~250 Token',
     RGBColor(0x6B, 0x3F, 0xC8), True, 'Identisch bei jeder\nAnfrage -> immer gecacht'),
    ('Kontext-Block\n(Top-5 Manual-Seiten)\n~2.000-6.000 Token',
     RGBColor(0x3A, 0x6B, 0xC8), True, 'Gleiche Seiten -> gleicher Block\nAufruf 2 teilt denselben Cache'),
    ('Nutzerfrage\n~10-30 Token',
     STEEL, False, 'Einmalig pro Anfrage\nNicht gecacht'),
]
bw2, bh2, by2 = Inches(3.1), Inches(2.2), Inches(1.8)
for j, (label, col, cached, note) in enumerate(blocks):
    bx2 = Inches(0.5) + j * Inches(4.2)
    rect(sl, bx2, by2, bw2, bh2, fill=col)
    txt(sl, label, bx2+Inches(0.15), by2+Inches(0.18), bw2-Inches(0.3), Inches(0.9),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tf, tt, tc = (YELLOW, 'GECACHT', DARK) if cached else (STEEL, 'NICHT GECACHT', MUTED)
    rect(sl, bx2+Inches(0.15), by2+Inches(1.1), bw2-Inches(0.3), Inches(0.45), fill=tf)
    txt(sl, tt, bx2+Inches(0.15), by2+Inches(1.12), bw2-Inches(0.3), Inches(0.35),
        size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(sl, note, bx2, by2+bh2+Inches(0.1), bw2, Inches(0.6),
        size=10, color=MUTED, align=PP_ALIGN.LEFT)

rows = [
    ('Block', 'Gecacht?', 'Warum'),
    ('System-Prompt', 'Ja', 'Identisch bei jeder Anfrage'),
    ('Kontext (Manual-Seiten)', 'Ja', 'Gleiche Seiten -> gleicher Block; Aufruf 2 nutzt Aufruf-1-Cache'),
    ('Nutzerfrage', 'Nein', 'Jede Frage einmalig'),
    ('Antwort (in Aufruf 2)', 'Nein', 'Jede Antwort einmalig'),
]
tw2 = [Inches(3.2), Inches(1.5), Inches(7.0)]
tx2 = [Inches(0.4), Inches(3.6), Inches(5.1)]
ty2, rh2 = Inches(4.5), Inches(0.38)
for ri, row in enumerate(rows):
    bg = STEEL if ri == 0 else (RGBColor(0x25, 0x2C, 0x42) if ri % 2 == 0
                                else RGBColor(0x1E, 0x24, 0x38))
    for ci, (cell, cw2, cx2) in enumerate(zip(row, tw2, tx2)):
        rect(sl, cx2, ty2+ri*rh2, cw2, rh2, fill=bg, line=STEEL, line_w=Pt(0.5))
        c2 = YELLOW if ri == 0 else (SUCCESS if (ci == 1 and 'Ja' in cell) else WHITE)
        txt(sl, cell, cx2+Inches(0.08), ty2+ri*rh2+Inches(0.05),
            cw2-Inches(0.1), rh2-Inches(0.05), size=10, bold=(ri == 0),
            color=c2, align=PP_ALIGN.LEFT)

txt(sl, ('Ersparnis: Gecachte Token kosten ~10x weniger. '
         'Bei Wiederholungsfragen sinken API-Kosten deutlich.'),
    Inches(0.4), Inches(6.53), Inches(12.2), Inches(0.45),
    size=11, color=MUTED, align=PP_ALIGN.LEFT)
footer(sl, 'Folie 5 / 7')


# -- Folie 6: Fehlercode-Lookup --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LIGHT)
header(sl, 'Stoerungs-Lookup -- ohne KI-Aufruf',
       'POST /errorcode  --  unter 100 ms, kein API-Aufruf erforderlich')

steps2 = [
    ('1', 'Code oder Stichwort eingeben',
     'Bediener gibt z. B. "1042" oder "Hydraulik" in das Suchfeld ein.'),
    ('2', 'Dict-Lookup im Arbeitsspeicher',
     'Exakter Code wird in errorcodes.json nachgeschlagen. Bei Stichwort: Keyword-Suche in allen Codes.'),
    ('3', 'Kein API-Aufruf',
     'Die Anthropic API wird nicht kontaktiert. Antwortzeit: unter 100 ms.'),
    ('4', 'Strukturiertes Ergebnis',
     'Fehlercode, Beschreibung, Ursache, Massnahme. Dazu Quellenlinks auf passende Manual-Seiten (top-3).'),
    ('5', 'Nachfrage zur KI weitergeben',
     'Inline-Eingabe direkt unter dem Ergebnis. Erst hier wird /ask ausgeloest.'),
]
for i, (num, title, detail) in enumerate(steps2):
    t = Inches(1.38) + i * Inches(1.0)
    rect(sl, Inches(0.3), t, Inches(0.5), Inches(0.5), fill=YELLOW)
    txt(sl, num, Inches(0.3), t+Inches(0.05), Inches(0.5), Inches(0.4),
        size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.0), t+Inches(0.04), Inches(8.0), Inches(0.38),
        size=13, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    txt(sl, detail, Inches(1.0), t+Inches(0.42), Inches(8.0), Inches(0.52),
        size=10, color=STEEL, align=PP_ALIGN.LEFT)

rect(sl, Inches(9.5), Inches(1.5), Inches(3.5), Inches(2.8),
     fill=DARK, line=YELLOW, line_w=Pt(2))
txt(sl, 'Zeitvergleich', Inches(9.65), Inches(1.6), Inches(3.2), Inches(0.4),
    size=13, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
for k, (label, val, col) in enumerate([
    ('Stoerungs-Lookup', '< 100 ms', SUCCESS),
    ('KI-Frage (/ask)',  '1-2 s',    WARNING),
    ('Seitenladung',     '0 ms',     SUCCESS),
]):
    ty3 = Inches(2.1) + k * Inches(0.65)
    txt(sl, label, Inches(9.65), ty3, Inches(2.0), Inches(0.4),
        size=11, color=WHITE, align=PP_ALIGN.LEFT)
    txt(sl, val, Inches(11.6), ty3, Inches(1.3), Inches(0.4),
        size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)
footer(sl, 'Folie 6 / 7')


# -- Folie 7: Deployment & naechste Schritte --
sl = add_slide()
rect(sl, 0, 0, W, H, fill=DARK)
header(sl, 'Deployment & offene Punkte',
       'Aktueller Stand und naechste Entwicklungsschritte')

rect(sl, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.2),
     fill=STEEL, line=YELLOW, line_w=Pt(1.5))
txt(sl, 'Aktuell in Betrieb', Inches(0.5), Inches(1.4), Inches(5.4), Inches(0.45),
    size=15, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
deploy = [
    ('Plattform',    'Railway  (Cloud PaaS)'),
    ('Auto-Deploy',  'Push auf main -> automatisch'),
    ('URL',          'aiassistant-production-da1f.up.railway.app'),
    ('Frontend',     'Single-File HTML  |  kein Build-Schritt'),
    ('Manual',       '2.180 HTML-Seiten im Repository'),
    ('Modell',       'claude-haiku-4-5-20251001'),
    ('PWA',          'Installierbar via Safari -> Home-Bildschirm'),
    ('Responsive',   'Desktop + Tablet/iPad  (Bottom-Nav <= 768 px)'),
    ('Mikrofon',     'Web Speech API  |  Audio bleibt auf dem Geraet'),
]
dy = Inches(1.93)
for label, value in deploy:
    txt(sl, label+':', Inches(0.5), dy, Inches(1.5), Inches(0.38),
        size=10, bold=True, color=MUTED, align=PP_ALIGN.LEFT)
    txt(sl, value, Inches(2.05), dy, Inches(3.9), Inches(0.38),
        size=10, color=WHITE, align=PP_ALIGN.LEFT)
    dy += Inches(0.46)

rect(sl, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.2),
     fill=RGBColor(0x1E, 0x24, 0x38), line=STEEL, line_w=Pt(1.5))
txt(sl, 'Offene Punkte (priorisiert)', Inches(6.7), Inches(1.4),
    Inches(6.1), Inches(0.45), size=15, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
next_steps = [
    ('1', 'Manual-Bilder extern hosten (240 MB -> Cloudflare R2)', WARNING),
    ('2', 'Authentifizierung (API-Key oder Liebherr AD/EntraID)',   WARNING),
    ('3', '/admin/reload  --  Manual-Aktualisierung ohne Neustart', MUTED),
    ('4', 'Gespraechsverlauf (letzte 2 Fragen als Kontext)',        MUTED),
    ('5', 'Branded PWA-Icons (192 + 512 px)',                       MUTED),
    ('6', 'Multi-Maschinen-Namespace (mehrere Manuale)',            MUTED),
]
ny = Inches(1.93)
for num, text, col in next_steps:
    rect(sl, Inches(6.65), ny+Inches(0.06), Inches(0.35), Inches(0.35), fill=col)
    txt(sl, num, Inches(6.65), ny+Inches(0.04), Inches(0.35), Inches(0.35),
        size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, text, Inches(7.1), ny+Inches(0.05), Inches(5.7), Inches(0.38),
        size=11, color=WHITE, align=PP_ALIGN.LEFT)
    ny += Inches(0.62)
footer(sl, 'Folie 7 / 7')


# -- Speichern --
out = 'KI_Assistent_Architektur.pptx'
prs.save(out)
print(f'Gespeichert: {out}')
